"""
示例：使用 python-pptx 创建一页 PPT
主题：大模型微调流程
布局：左侧 4 个竖向流程方框 + 箭头连接 + 右侧对应说明
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml import parse_xml


# ── 颜色常量 ──────────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1F, 0x1F, 0x1F)
RED_LINE = RGBColor(0xC0, 0x00, 0x00)
BOX_FILL = RGBColor(0xE8, 0xF3, 0xFF)
BOX_BORDER = RGBColor(0x5B, 0x9B, 0xD5)
BOX_TEXT = RGBColor(0x1A, 0x56, 0x9A)
DESC_TEXT = RGBColor(0x33, 0x33, 0x33)
ARROW_COLOR = RGBColor(0x5B, 0x9B, 0xD5)

# ── 字体与语言 ────────────────────────────────────────────────────────────────
FONT = "微软雅黑"
# 必须为每个 run 设置 language_id，否则 PowerPoint 会按英文进行拼写校对，
# 导致中文文本出现大量红色波浪线（误报为拼写错误）。
LANG_ZH = MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE

# ── 幻灯片尺寸（16:9 宽屏）──────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.33, 7.5

# ── 微调流程内容 ───────────────────────────────────────────────────────────────
STEPS = [
    {
        "title": "1. 准备微调的数据集",
        "desc": (
            "收集并清洗与目标场景相关的数据，统一为指令-回答、"
            "问答对或分类样本等格式，并划分训练集、验证集与测试集。"
        ),
    },
    {
        "title": "2. 选择基座模型",
        "desc": (
            "根据任务复杂度、参数规模、推理成本和开源许可证，"
            "选择合适的基础模型作为微调起点。"
        ),
    },
    {
        "title": "3. 选择微调框架",
        "desc": (
            "结合资源条件选择 LoRA、QLoRA、全参数微调或 SFT 框架，"
            "并配置训练超参数、显存策略与评估流程。"
        ),
    },
    {
        "title": "4. 测试微调效果",
        "desc": (
            "同时评估专用任务效果与通用任务能力，关注准确率、"
            "鲁棒性、幻觉率以及是否出现泛化能力下降。"
        ),
    },
]


def set_white_background(slide):
    """将幻灯片背景设置为纯白色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title(slide, text):
    """添加居中标题"""
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.18),
        Inches(12.33), Inches(0.75),
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.language_id = LANG_ZH


def add_red_line(slide):
    """在标题下方添加分隔线"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.6), Inches(1.0),
        Inches(12.73), Inches(1.0),
    )
    line.line.color.rgb = RED_LINE
    line.line.width = Pt(2.0)


def add_flow_box(slide, left, top, width, height, text):
    """添加左侧圆角流程方框"""
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BOX_FILL
    box.line.color.rgb = BOX_BORDER
    box.line.width = Pt(1.8)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BOX_TEXT
    run.font.language_id = LANG_ZH

    return box


def add_description(slide, left, top, width, height, text):
    """添加右侧流程说明文字"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.color.rgb = DESC_TEXT
    run.font.language_id = LANG_ZH


def add_down_connector(slide, upper_box, lower_box):
    """将连接符锚定到上下两个流程方框"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        0,
        0,
        0,
        0,
    )
    connector.begin_connect(upper_box, 2)  # 下中
    connector.end_connect(lower_box, 0)    # 上中
    connector.line.color.rgb = ARROW_COLOR
    connector.line.width = Pt(1.4)

    ln = connector.line._get_or_add_ln()
    ln.append(
        parse_xml(
            '<a:tailEnd type="triangle" w="lg" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(slide)
    add_title(slide, "大模型微调流程")
    add_red_line(slide)

    flow_left = 0.7
    flow_top = 1.35
    flow_w = 3.2
    flow_h = 1.05
    flow_gap = 0.28

    desc_left = 4.35
    desc_w = 8.1
    desc_h = flow_h

    prev_box = None

    for index, step in enumerate(STEPS):
        top = flow_top + index * (flow_h + flow_gap)

        flow_box = add_flow_box(
            slide,
            Inches(flow_left),
            Inches(top),
            Inches(flow_w),
            Inches(flow_h),
            step["title"],
        )

        add_description(
            slide,
            Inches(desc_left),
            Inches(top + 0.02),
            Inches(desc_w),
            Inches(desc_h),
            step["desc"],
        )

        if prev_box is not None:
            add_down_connector(slide, prev_box, flow_box)

        prev_box = flow_box

    output_path = "tests/output-example4-finetune-flow.pptx"
    prs.save(output_path)
    print(f"已保存：{output_path}")


if __name__ == "__main__":
    main()
