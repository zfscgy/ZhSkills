"""
示例：使用 python-pptx 创建一页 PPT
主题：大模型训练的三个阶段（预训练 → 微调 → RLHF）
布局：标题 + 红色分隔线 + 3 张卡片水平排列（蓝色右箭头连接）
卡片结构：上半部分为图片占位矩形，下半部分为阶段标题与文字说明
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 颜色常量 ──────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_LINE    = RGBColor(0xC0, 0x00, 0x00)   # 深红分隔线
GRAY_BG     = RGBColor(0xF2, 0xF2, 0xF2)   # 卡片浅灰背景
TITLE_COLOR = RGBColor(0x1F, 0x1F, 0x1F)   # 主标题近黑色
CARD_TITLE  = RGBColor(0x1A, 0x56, 0x9A)   # 卡片标题蓝色
CARD_TEXT   = RGBColor(0x33, 0x33, 0x33)   # 卡片正文深灰
IMG_BG      = RGBColor(0xD6, 0xE4, 0xF0)   # 图片占位区浅蓝灰
IMG_BORDER  = RGBColor(0x9B, 0xC2, 0xE6)   # 占位区边框蓝
IMG_TEXT    = RGBColor(0x70, 0x8E, 0xA8)   # 占位符提示文字
ARROW_COLOR = RGBColor(0x1A, 0x56, 0x9A)   # 箭头蓝色（与卡片标题同色系）

# ── 字体 ─────────────────────────────────────────────────────────────────────
FONT = "微软雅黑"

# ── 幻灯片尺寸（16:9 宽屏）──────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.33, 7.5

# ── 三个训练阶段的内容 ─────────────────────────────────────────────────────────
STAGES = [
    {
        "title": "① 预训练（Pre-training）",
        "content": (
            "在万亿 Token 级无标注语料上，以自回归/掩码语言模型为目标"
            "训练 Transformer，使模型习得语言规律、世界知识与推理能力，"
            "形成通用基础模型（Base Model）。"
            "代表：GPT-4、LLaMA-3、Qwen。"
        ),
    },
    {
        "title": "② 监督微调（SFT）",
        "content": (
            "使用人工标注的高质量「指令—回答」对数据对基础模型进行"
            "有监督微调，引导模型学习特定回答风格，从「文本补全」"
            "转变为「遵循指令」，大幅提升实用性与可控性。"
            "代表：InstructGPT、Alpaca、Vicuna。"
        ),
    },
    {
        "title": "③ RLHF（人类反馈强化学习）",
        "content": (
            "训练奖励模型（Reward Model）对回答质量打分，"
            "再以 PPO 算法对 SFT 模型做强化学习优化，"
            "使输出更符合人类偏好与安全要求。"
            "DPO 等变体进一步简化了训练流程，是当前主流对齐技术。"
        ),
    },
]


# ─────────────────────────────────────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────────────────────────────────────

def set_white_background(slide):
    """将幻灯片背景设置为纯白色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title(slide, text):
    """添加居中大标题"""
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(12.33), Inches(0.9),
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_red_line(slide):
    """在标题下方添加红色水平分隔线"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.5),   Inches(1.15),
        Inches(12.83), Inches(1.15),
    )
    line.line.color.rgb = RED_LINE
    line.line.width = Pt(2.5)


def add_image_placeholder(slide, left, top, width, height):
    """
    添加图片占位矩形：浅蓝背景 + 蓝色边框 + 居中提示文字。
    用户可在 PPT 中删除此形状后插入实际图片。
    """
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = IMG_BG
    rect.line.color.rgb = IMG_BORDER
    rect.line.width = Pt(1.5)

    tf = rect.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # 文字垂直居中

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[ 请在此处插入图片 ]"
    run.font.name = FONT
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = IMG_TEXT


def add_card(slide, left, top, width, height, title_text, body_text):
    """
    添加三段式卡片：
      - 外层浅灰矩形（整张卡片）
      - 上部：图片占位矩形（约占 48% 高度）
      - 下部：阶段标题 + 正文文本框
    """
    IMG_RATIO = 0.48    # 图片占位区占卡片高度的比例
    PAD = Inches(0.18)  # 内边距

    # 外层卡片背景矩形（shape type 1 = rectangle）
    card_bg = slide.shapes.add_shape(1, left, top, width, height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = GRAY_BG
    card_bg.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    card_bg.line.width = Pt(0.75)

    # 图片占位区
    img_h = height * IMG_RATIO
    add_image_placeholder(
        slide,
        left + PAD,
        top + PAD,
        width - PAD * 2,
        img_h - PAD,
    )

    # 文字区域（卡片标题 + 正文）
    txt_top = top + img_h + PAD * 0.5
    txt_h   = height - img_h - PAD * 1.5
    txBox = slide.shapes.add_textbox(
        left + PAD,
        txt_top,
        width - PAD * 2,
        txt_h,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # 卡片阶段标题
    p_title = tf.paragraphs[0]
    p_title.space_after = Pt(5)
    run_t = p_title.add_run()
    run_t.text = title_text
    run_t.font.name = FONT
    run_t.font.size = Pt(17)
    run_t.font.bold = True
    run_t.font.color.rgb = CARD_TITLE

    # 正文段落
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(2)
    run_b = p_body.add_run()
    run_b.text = body_text
    run_b.font.name = FONT
    run_b.font.size = Pt(12)
    run_b.font.color.rgb = CARD_TEXT


def add_right_arrow(slide, left, top, width, height):
    """
    在两张卡片之间添加实心右箭头。
    MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW = 33（可通过枚举验证）
    """
    arrow = slide.shapes.add_shape(33, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ARROW_COLOR
    # 边框颜色与填充一致，视觉上无边框
    arrow.line.color.rgb = ARROW_COLOR


# ─────────────────────────────────────────────────────────────────────────────
# 主函数
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    set_white_background(slide)
    add_title(slide, "大模型训练的三个阶段")
    add_red_line(slide)

    # ── 布局参数（单位 Inches）────────────────────────────────────────────────
    MARGIN_X = 0.4    # 左右外边距
    MARGIN_Y = 1.3    # 卡片顶部起始位置（红线下方）
    ARROW_W  = 0.55   # 箭头宽度
    ARROW_H  = 0.4    # 箭头高度
    GAP      = 0.12   # 卡片与箭头之间的空隙

    # 3 张卡片 + 2 个箭头间隔，均分可用宽度
    total_gap = (ARROW_W + GAP * 2) * 2
    CARD_W = (SLIDE_W - MARGIN_X * 2 - total_gap) / 3   # ≈ 3.76"
    CARD_H = SLIDE_H - MARGIN_Y - 0.2                   # ≈ 6.0"

    for i, stage in enumerate(STAGES):
        card_left = MARGIN_X + i * (CARD_W + GAP + ARROW_W + GAP)

        add_card(
            slide,
            Inches(card_left),
            Inches(MARGIN_Y),
            Inches(CARD_W),
            Inches(CARD_H),
            stage["title"],
            stage["content"],
        )

        # 在卡片 0→1 和 1→2 之间各放一个右箭头，垂直居中于图片占位区
        if i < 2:
            arrow_left = card_left + CARD_W + GAP
            # 箭头对齐图片占位区的垂直中心（图片区占卡片 48%）
            img_center_y = MARGIN_Y + CARD_H * 0.48 * 0.5
            arrow_top = img_center_y - ARROW_H / 2
            add_right_arrow(
                slide,
                Inches(arrow_left),
                Inches(arrow_top),
                Inches(ARROW_W),
                Inches(ARROW_H),
            )

    output_path = "tests/output-example2.pptx"
    prs.save(output_path)
    print(f"已保存：{output_path}")


if __name__ == "__main__":
    main()
