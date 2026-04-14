"""
示例：使用 python-pptx 创建一页 PPT
主题：三大顶尖大模型厂商介绍（OpenAI / Anthropic / Google）
布局：标题 + 红色分隔线 + 3 行（每行：左侧图片占位符 + 右侧文字）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 颜色常量 ──────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_LINE    = RGBColor(0xC0, 0x00, 0x00)   # 深红分隔线
GRAY_BG     = RGBColor(0xF2, 0xF2, 0xF2)   # 行卡片浅灰背景
TITLE_COLOR = RGBColor(0x1F, 0x1F, 0x1F)   # 主标题近黑色
CARD_TITLE  = RGBColor(0x1A, 0x56, 0x9A)   # 公司名称蓝色
CARD_TEXT   = RGBColor(0x33, 0x33, 0x33)   # 正文深灰
IMG_BG      = RGBColor(0xD6, 0xE4, 0xF0)   # 图片占位区浅蓝灰
IMG_BORDER  = RGBColor(0x9B, 0xC2, 0xE6)   # 占位区边框蓝
IMG_TEXT    = RGBColor(0x70, 0x8E, 0xA8)   # 占位符提示文字

# ── 字体 ─────────────────────────────────────────────────────────────────────
FONT = "微软雅黑"

# ── 幻灯片尺寸（16:9 宽屏）──────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.33, 7.5

# ── 三家公司的内容（数据截至 2026 年 3 月）─────────────────────────────────────
COMPANIES = [
    {
        "title": "OpenAI",
        "subtitle": "GPT-5.4 · o3/o4-mini · 估值 8400 亿美元 · 年化营收 250 亿（2026.03）",
        "content": (
            "2026 年 2 月完成史上最大私募融资：亚马逊、英伟达、软银合投 1100 亿美元，"
            "3 月追加 100 亿，估值达 8400 亿美元，剑指万亿 IPO。"
            "年化营收突破 250 亿美元，ChatGPT 周活跃用户达 9 亿。"
            "3 月发布旗舰模型 GPT-5.4，深度融合推理、编程与自主 Agent 操作能力，"
            "是当前最广泛商业落地的大模型平台。"
        ),
    },
    {
        "title": "Anthropic",
        "subtitle": "Claude Opus/Sonnet 4.6 · 年化营收 140 亿 · 估值 3800 亿（2026.02）",
        "content": (
            "2026 年 2 月完成 G 轮 300 亿美元融资，估值 3800 亿，跃居全球第二大私人科技公司。"
            "年化营收达 140 亿美元，仅 Claude Code 一款产品年化营收即超 25 亿。"
            "Opus 4.6 与 Sonnet 4.6 相继发布，支持混合推理（即时 / 深度思考双模式）；"
            "宪法 AI 对齐方法持续引领安全研究，Claude 5 正在研发中。"
        ),
    },
    {
        "title": "Google DeepMind",
        "subtitle": "Gemini 3.1 Pro · ARC-AGI-2 77.1% · 多模态 Embedding（2026.02）",
        "content": (
            "2026 年 2 月发布 Gemini 3.1 Pro，在 Humanity's Last Exam 得分 44.4%，"
            "ARC-AGI-2 逻辑推理达 77.1%（前代仅 31.1%），复杂问题求解能力大幅跃升。"
            "同步推出全球首个原生多模态 Embedding 模型 Gemini Embedding 2，"
            "统一处理文本、图像、视频、音频与文档。Gemini 深度整合 Google 搜索、"
            "Vertex AI 与 NotebookLM，覆盖消费与企业双侧生态。"
        ),
    },
]


# ─────────────────────────────────────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────────────────────────────────────

def set_white_background(slide):
    """将幻灯片背景设置为纯白色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title(slide, text):
    """添加居中大标题"""
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(12.33), Inches(0.9),
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_red_line(slide):
    """在标题下方添加红色水平分隔线"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.5),   Inches(1.12),
        Inches(12.83), Inches(1.12),
    )
    line.line.color.rgb = RED_LINE
    line.line.width = Pt(2.5)


def add_image_placeholder(slide, left, top, width, height):
    """
    图片占位矩形：浅蓝背景 + 无边框无阴影 + 垂直居中提示文字。
    用户可删除此形状后粘贴公司 Logo 或相关图片。
    """
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = IMG_BG
    rect.line.fill.background()
    rect.shadow.inherit = False

    tf = rect.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[ 请插入图片 ]"
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = IMG_TEXT


def add_row(slide, left, top, width, height,
            title_text, subtitle_text, body_text):
    """
    添加一行卡片：
      - 整行浅灰背景矩形
      - 左侧：正方形图片占位符
      - 右侧：公司名称（大）+ 副标题（中）+ 正文（小）
    """
    PAD   = Inches(0.18)   # 内边距
    IMG_W = height - PAD   # 图片占位宽度 ≈ 行高，近似正方形
    GAP_X = Inches(0.22)   # 图片与文字之间的间隙

    # 整行背景矩形
    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRAY_BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    # 左侧图片占位符（近正方形）
    add_image_placeholder(
        slide,
        left + PAD,
        top + PAD * 0.5,
        IMG_W,
        height - PAD,
    )

    # 右侧文字区域
    txt_left = left + PAD + IMG_W + GAP_X
    txt_top  = top + PAD * 0.6
    txt_w    = width - PAD - IMG_W - GAP_X - PAD
    txt_h    = height - PAD * 1.2

    txBox = slide.shapes.add_textbox(txt_left, txt_top, txt_w, txt_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # 公司名称（最大）
    p_name = tf.paragraphs[0]
    p_name.space_after = Pt(2)
    r = p_name.add_run()
    r.text = title_text
    r.font.name = FONT
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = CARD_TITLE

    # 副标题（产品/关键词标签，中等字号）
    p_sub = tf.add_paragraph()
    p_sub.space_before = Pt(1)
    p_sub.space_after  = Pt(4)
    r_sub = p_sub.add_run()
    r_sub.text = subtitle_text
    r_sub.font.name = FONT
    r_sub.font.size = Pt(16)
    r_sub.font.bold = False
    r_sub.font.color.rgb = RGBColor(0x70, 0x70, 0x70)   # 中灰，区别于正文

    # 正文介绍
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(1)
    r_body = p_body.add_run()
    r_body.text = body_text
    r_body.font.name = FONT
    r_body.font.size = Pt(14)
    r_body.font.color.rgb = CARD_TEXT


# ─────────────────────────────────────────────────────────────────────────────
# 主函数
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    set_white_background(slide)
    add_title(slide, "三大顶尖大模型厂商")
    add_red_line(slide)

    # ── 布局参数（单位 Inches）────────────────────────────────────────────────
    MARGIN_X = 0.4    # 左右外边距
    MARGIN_Y = 1.27   # 第一行顶部起始（红线下方留白）
    GAP_Y    = 0.18   # 行间距
    ROW_W    = SLIDE_W - MARGIN_X * 2          # 每行宽度 = 12.53"
    # 3 行均分剩余高度
    ROW_H    = (SLIDE_H - MARGIN_Y - 0.12 - GAP_Y * 2) / 3   # ≈ 1.93"

    for i, company in enumerate(COMPANIES):
        row_top = MARGIN_Y + i * (ROW_H + GAP_Y)
        add_row(
            slide,
            Inches(MARGIN_X),
            Inches(row_top),
            Inches(ROW_W),
            Inches(ROW_H),
            company["title"],
            company["subtitle"],
            company["content"],
        )

    output_path = "tests/output-example3.pptx"
    prs.save(output_path)
    print(f"已保存：{output_path}")


if __name__ == "__main__":
    main()
