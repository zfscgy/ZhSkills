"""
示例：使用 python-pptx 创建一页 PPT
主题：大模型发展的四个阶段
布局：标题 + 红色分隔线 + 2x2 文本块
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.util import Inches, Pt

# ── 颜色常量 ──────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_LINE    = RGBColor(0xC0, 0x00, 0x00)   # 深红，视觉上比纯红更专业
GRAY_BG     = RGBColor(0xF2, 0xF2, 0xF2)   # 浅灰背景
TITLE_COLOR = RGBColor(0x1F, 0x1F, 0x1F)   # 近黑色标题
CARD_TITLE  = RGBColor(0x1A, 0x56, 0x9A)   # 蓝色卡片标题
CARD_TEXT   = RGBColor(0x33, 0x33, 0x33)   # 深灰正文

# ── 字体与语言 ────────────────────────────────────────────────────────────────
FONT = "微软雅黑"
# 必须为每个 run 设置 language_id，否则 PowerPoint 会按英文进行拼写校对，
# 导致中文文本出现大量红色波浪线（误报为拼写错误）。
LANG_ZH = MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE

# ── 幻灯片尺寸（标准 16:9，单位 Inches）─────────────────────────────────────
SLIDE_W, SLIDE_H = 13.33, 7.5

# ── 四个阶段的内容 ─────────────────────────────────────────────────────────────
CARDS = [
    {
        "title": "第一阶段：语言模型奠基期",
        "content": (
            "以 GPT-2、BERT 为代表，Transformer 架构确立主导地位。"
            "研究重心在于自监督预训练范式，通过海量无标注文本学习通用语言表征，"
            "验证了规模扩展的可行性，为后续千亿参数模型铺平道路。"
        ),
    },
    {
        "title": "第二阶段：规模化涌现期",
        "content": (
            "GPT-3（1750亿参数）展示出「涌现能力」：少样本学习、链式推理等能力"
            "随参数量增长而突然出现。PaLM、Chinchilla 等工作进一步揭示"
            "数据量与模型规模的最优配比，奠定了大模型的工程化基础。"
        ),
    },
    {
        "title": "第三阶段：对齐与指令微调期",
        "content": (
            "InstructGPT 引入 RLHF（人类反馈强化学习），使模型从「补全文本」"
            "转变为「遵循指令」。ChatGPT 的爆发式普及证明对齐技术的商业价值，"
            "Llama 系列开源推动社区繁荣，SFT + RLHF 成为行业标配流程。"
        ),
    },
    {
        "title": "第四阶段：多模态与智能体期",
        "content": (
            "GPT-4V、Gemini 将视觉、音频融入统一模型；Function Calling、"
            "RAG 赋予模型外部工具调用能力。以 AutoGPT、Claude Computer Use"
            "为代表的 Agent 框架，使大模型具备规划、记忆与自主执行复杂任务的能力。"
        ),
    },
]


def set_white_background(slide):
    """将幻灯片背景设置为纯白色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title(slide):
    """添加居中标题文本框"""
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.25),      # left, top
        Inches(11.33), Inches(0.9),     # width, height
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "大模型发展的四个阶段"
    run.font.name = FONT
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.language_id = LANG_ZH


def add_red_line(slide):
    """在标题下方添加一条红色水平分隔线"""
    # add_connector(connector_type, begin_x, begin_y, end_x, end_y)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.5),  Inches(1.25),    # 起点 (x1, y1)
        Inches(12.83), Inches(1.25),   # 终点 (x2, y2)
    )
    line.line.color.rgb = RED_LINE
    line.line.width = Pt(2.5)


def add_card(slide, left, top, width, height, title_text, body_text):
    """
    添加一个带浅灰背景的文本卡片。
    先用矩形做背景色，再叠加透明文本框写内容，
    这样可以精确控制背景与文字的独立格式。
    """
    # ── 背景矩形 ──────────────────────────────────────────────────────────────
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt as PtU
    rect = slide.shapes.add_shape(
        1,              # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = GRAY_BG
    rect.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)   # 浅灰边框
    rect.line.width = Pt(0.75)

    # ── 文本框（覆盖在矩形上，透明背景）────────────────────────────────────────
    PADDING = Inches(0.2)
    txBox = slide.shapes.add_textbox(
        left + PADDING,
        top + PADDING,
        width - PADDING * 2,
        height - PADDING * 2,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # 卡片标题段落
    p_title = tf.paragraphs[0]
    p_title.space_after = Pt(6)
    run_title = p_title.add_run()
    run_title.text = title_text
    run_title.font.name = FONT
    run_title.font.size = Pt(18)
    run_title.font.bold = True
    run_title.font.color.rgb = CARD_TITLE
    run_title.font.language_id = LANG_ZH

    # 正文段落
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(4)
    run_body = p_body.add_run()
    run_body.text = body_text
    run_body.font.name = FONT
    run_body.font.size = Pt(14)
    run_body.font.color.rgb = CARD_TEXT
    run_body.font.language_id = LANG_ZH


def main():
    prs = Presentation()

    # 设置幻灯片尺寸为 16:9 宽屏
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 使用空白版式（index=6）
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ── 背景 / 标题 / 分隔线 ──────────────────────────────────────────────────
    set_white_background(slide)
    add_title(slide)
    add_red_line(slide)

    # ── 四个卡片：2 列 × 2 行 ─────────────────────────────────────────────────
    # 布局参数（单位 Inches）
    MARGIN_X = 0.4          # 左右外边距
    MARGIN_Y = 1.4          # 顶部起始位置（红线下方）
    GAP_X    = 0.25         # 列间距
    GAP_Y    = 0.25         # 行间距
    CARD_W   = (SLIDE_W - MARGIN_X * 2 - GAP_X) / 2   # 每张卡片宽度 ≈ 6.12"
    CARD_H   = (SLIDE_H - MARGIN_Y - 0.2 - GAP_Y) / 2  # 每张卡片高度 ≈ 2.83"

    positions = [
        (0, 0),  # 左上
        (1, 0),  # 右上
        (0, 1),  # 左下
        (1, 1),  # 右下
    ]

    for i, (col, row) in enumerate(positions):
        left = Inches(MARGIN_X + col * (CARD_W + GAP_X))
        top  = Inches(MARGIN_Y + row * (CARD_H + GAP_Y))
        add_card(
            slide,
            left, top,
            Inches(CARD_W), Inches(CARD_H),
            CARDS[i]["title"],
            CARDS[i]["content"],
        )

    # ── 保存 ──────────────────────────────────────────────────────────────────
    output_path = "tests/output-example1.pptx"
    prs.save(output_path)
    print(f"已保存：{output_path}")


if __name__ == "__main__":
    main()
