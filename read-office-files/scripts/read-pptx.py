"""
read-pptx.py  —  读取 .pptx 文件，输出结构化文本
用法:
    python read-pptx.py <file.pptx> [--json] [--notes] [--tables] [--meta]

选项:
    --json     以 JSON 格式输出（便于程序处理）
    --notes    同时提取演讲者备注
    --tables   同时提取表格内容
    --meta     同时输出演示文稿元数据

依赖: pip install python-pptx
"""

import sys
import json
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER


def is_title_placeholder(shape) -> bool:
    """判断占位符是否为标题类型。"""
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    return ph_type in (
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
    )


def extract_shape_text(shape) -> list[str]:
    """从单个 shape 中提取所有非空文本段落。"""
    if not shape.has_text_frame:
        return []
    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    return lines


def extract_table_from_shape(shape) -> list[list[str]] | None:
    """从表格 shape 中提取内容，非表格返回 None。"""
    if not shape.has_table:
        return None
    rows = []
    for row in shape.table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows


def extract_slide(slide, slide_num: int, include_notes: bool, include_tables: bool) -> dict:
    """提取单张幻灯片的所有内容。"""
    title = ""
    body_lines = []
    tables = []

    for shape in slide.shapes:
        # 提取表格
        if include_tables and shape.has_table:
            tbl = extract_table_from_shape(shape)
            if tbl:
                tables.append(tbl)
            continue

        if not shape.has_text_frame:
            continue

        if is_title_placeholder(shape):
            title = shape.text_frame.text.strip()
        else:
            body_lines.extend(extract_shape_text(shape))

    # 演讲者备注
    notes_text = ""
    if include_notes and slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf:
            notes_text = notes_tf.text.strip()

    result = {
        "slide": slide_num,
        "title": title,
        "body": body_lines,
    }
    if include_tables and tables:
        result["tables"] = tables
    if include_notes:
        result["notes"] = notes_text
    return result


def extract_meta(prs) -> dict:
    """提取演示文稿核心属性。"""
    cp = prs.core_properties
    return {
        "title":    cp.title or "",
        "author":   cp.author or "",
        "subject":  cp.subject or "",
        "keywords": cp.keywords or "",
        "created":  str(cp.created) if cp.created else "",
        "modified": str(cp.modified) if cp.modified else "",
        "slides":   len(prs.slides),
    }


def to_markdown(slides: list[dict]) -> str:
    """将幻灯片数据渲染为 Markdown 文本。"""
    lines = []
    for slide in slides:
        num = slide["slide"]
        title = slide.get("title", "")
        body = slide.get("body", [])
        tables = slide.get("tables", [])
        notes = slide.get("notes", "")

        lines.append(f"## Slide {num}" + (f": {title}" if title else ""))
        lines.append("")

        for line in body:
            lines.append(f"- {line}")
        if body:
            lines.append("")

        if tables:
            for i, table in enumerate(tables, 1):
                if not table:
                    continue
                lines.append(f"*表格 {i}*")
                header = table[0]
                lines.append("| " + " | ".join(header) + " |")
                lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in table[1:]:
                    lines.append("| " + " | ".join(row) + " |")
                lines.append("")

        if notes:
            lines.append(f"> 备注：{notes}")
            lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(description="读取 .pptx 文件内容")
    parser.add_argument("file", help=".pptx 文件路径")
    parser.add_argument("--json",   action="store_true", help="JSON 格式输出")
    parser.add_argument("--notes",  action="store_true", help="提取演讲者备注")
    parser.add_argument("--tables", action="store_true", help="提取表格")
    parser.add_argument("--meta",   action="store_true", help="输出元数据")
    args = parser.parse_args()

    path = Path(args.file)
    if not path.exists():
        print(f"[错误] 文件不存在：{path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(path))
    slides_data = [
        extract_slide(slide, i + 1, args.notes, args.tables)
        for i, slide in enumerate(prs.slides)
    ]
    meta = extract_meta(prs) if args.meta else None

    if args.json:
        output = {"slides": slides_data}
        if meta is not None:
            output["meta"] = meta
        print(json.dumps(output, ensure_ascii=False, indent=2))
    else:
        if meta:
            print("=== 演示文稿元数据 ===")
            for k, v in meta.items():
                if v:
                    print(f"{k}: {v}")
            print()
        print(to_markdown(slides_data))


if __name__ == "__main__":
    main()
